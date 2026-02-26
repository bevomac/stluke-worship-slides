"""Build .pptx from parsed sections using template layouts.

Maps each parsed section to the correct slide layout from the template and
populates placeholders with formatted text.
"""

import re
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import (
    GOLD, RED, WHITE, HYMN_NUM_GOLD, PEACH,
    TITLE_SIZE, SUBTITLE_SIZE, BODY_SIZE, BODY_LARGE, BODY_SMALL,
    READING_SIZE, HYMN_LYRIC_SIZE, VERSE_NUM_SIZE, PAGE_REF_SIZE,
    SMALL_TEXT, LICENSE_SIZE,
    TEMPLATE_PATH, LAYOUTS,
    SLIDE_WIDTH, SLIDE_HEIGHT,
)
from text_utils import (
    split_text_into_slides, split_dialogue_into_slides,
    split_reading_text, is_congregation_line, get_speaker_prefix,
)


class SlideGenerator:
    """Generates a PowerPoint presentation from parsed worship sections."""

    def __init__(self, template_path=None):
        self.template_path = template_path or TEMPLATE_PATH
        self.prs = Presentation(self.template_path)
        self._layout_cache = {}
        for i, layout in enumerate(self.prs.slide_layouts):
            self._layout_cache[layout.name] = layout

    def get_layout(self, name):
        """Get a slide layout by name."""
        if name in self._layout_cache:
            return self._layout_cache[name]
        raise KeyError(f"Layout '{name}' not found. Available: {list(self._layout_cache.keys())}")

    def add_slide(self, layout_name):
        """Add a slide with the given layout and return it."""
        layout = self.get_layout(layout_name)
        return self.prs.slides.add_slide(layout)

    def generate(self, sections):
        """Generate all slides from parsed sections.

        Args:
            sections: List of section dicts from docx_parser.parse_worship_script()

        Returns:
            The Presentation object.
        """
        for section in sections:
            handler = self._get_handler(section['type'])
            if handler:
                handler(section)

        return self.prs

    def save(self, output_path):
        """Save the presentation to a file."""
        self.prs.save(output_path)

    # ── Section handlers ────────────────────────────────────────────────

    def _get_handler(self, section_type):
        """Map section type to its handler method."""
        handlers = {
            'title': self._gen_title_slide,
            'announcements': self._gen_announcements,
            'prelude': self._gen_prelude,
            'confession': self._gen_confession,
            'gathering_hymn': self._gen_hymn,
            'greeting': self._gen_greeting,
            'kyrie': self._gen_kyrie,
            'canticle_of_praise': self._gen_kyrie,  # Same treatment - media slide
            'prayer_of_day': self._gen_prayer_of_day,
            'childrens_message': self._gen_childrens_message,
            'first_reading': self._gen_reading,
            'psalm': self._gen_reading,
            'second_reading': self._gen_reading,
            'gospel_acclamation': self._gen_gospel_acclamation,
            'gospel_announcement': self._gen_gospel_announcement,
            'gospel_reading': self._gen_gospel_reading,
            'message': self._gen_message,
            'hymn_of_day': self._gen_hymn,
            'creed': self._gen_creed,
            'prayers': self._gen_prayers,
            'peace': self._gen_peace,
            'offering': self._gen_offering,
            'offertory_hymn': self._gen_offertory,
            'offertory_prayer': self._gen_offertory_prayer,
            'great_thanksgiving': self._gen_great_thanksgiving,
            'words_of_institution': self._gen_words_of_institution,
            'holy_holy_holy': self._gen_media_placeholder,
            'lords_prayer': self._gen_lords_prayer,
            'communion_invitation': self._gen_communion_invitation,
            'communion_hymn': self._gen_communion_hymn,
            'lamb_of_god': self._gen_media_placeholder,
            'post_communion_blessing': self._gen_post_communion_blessing,
            'post_communion_canticle': self._gen_post_communion_canticle,
            'post_communion_prayer': self._gen_post_communion_prayer,
            'blessing': self._gen_blessing,
            'sending_hymn': self._gen_hymn,
            'dismissal': self._gen_dismissal,
            'postlude': self._gen_postlude,
            'generic_section': self._gen_generic_section,
        }
        return handlers.get(section_type)

    # ── Title Slide ─────────────────────────────────────────────────────

    def _gen_title_slide(self, section):
        """Generate the title slide with service name, date, pastor, etc."""
        slide = self.add_slide('Title Slide')
        meta = section.get('metadata', {})

        # Placeholder 11: Service name (e.g., "Worship Service of Holy Communion")
        ph11 = self._get_placeholder(slide, 11)
        if ph11:
            self._set_text(ph11, meta.get('service_name', 'Worship Service'),
                           size=Pt(39), bold=False)

        # Placeholder 12: Sunday/season name (e.g., "First Sunday\nin Lent")
        ph12 = self._get_placeholder(slide, 12)
        if ph12:
            sunday = meta.get('sunday_name', '')
            self._set_text(ph12, sunday)

        # Placeholder 13: Pastor and Music Director
        ph13 = self._get_placeholder(slide, 13)
        if ph13:
            tf = ph13.text_frame
            tf.clear()
            pastor = meta.get('pastor', '')
            music = meta.get('music_director', '')
            if pastor:
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = pastor
                run.font.bold = True
                run.font.size = Pt(35)
                run.font.color.rgb = PEACH
            if music:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = music
                run.font.bold = True
                run.font.size = Pt(35)
                run.font.color.rgb = PEACH

    # ── Announcements ───────────────────────────────────────────────────

    def _gen_announcements(self, section):
        """Generate announcements slide."""
        slide = self.add_slide('Title and Content')
        ph0 = self._get_placeholder(slide, 0)
        if ph0:
            self._set_text(ph0, 'Announcements', size=TITLE_SIZE)

    # ── Prelude ─────────────────────────────────────────────────────────

    def _gen_prelude(self, section):
        """Generate prelude slide."""
        slide = self.add_slide('Title and Content')
        ph0 = self._get_placeholder(slide, 0)
        if ph0:
            self._set_text(ph0, 'Prelude', size=TITLE_SIZE)

        # Add performer name if available
        content = section.get('content', [])
        if content:
            # Add a text box for the performer
            performer = '\n'.join(content)
            self._add_textbox(slide, performer,
                              left=Inches(1), top=Inches(4.5),
                              width=Inches(8), height=Inches(1),
                              size=BODY_SIZE)

    # ── Confession and Forgiveness ──────────────────────────────────────

    def _gen_confession(self, section):
        """Generate Confession and Forgiveness slides (2-4 slides)."""
        content = section.get('content', [])

        # First slide: section header with beginning of text
        chunks = split_text_into_slides('\n'.join(content), max_chars=400, max_lines=7)

        if not chunks:
            # Just the header
            slide = self.add_slide('4_Section Header')
            self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)
            return

        # First slide uses 4_Section Header
        slide = self.add_slide('4_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)
        self._format_dialogue_in_placeholder(slide, 1, chunks[0], size=Pt(39))

        # Overflow slides use 1_Title and Content
        for chunk in chunks[1:]:
            slide = self.add_slide('1_Title and Content')
            self._format_dialogue_in_placeholder(slide, 1, chunk, size=Pt(40))

    # ── Hymns (Gathering, Day, Sending) ─────────────────────────────────

    def _gen_hymn(self, section):
        """Generate hymn slides: intro slide + one slide per verse."""
        meta = section.get('metadata', {})
        hymn_title = meta.get('hymn_title', '')
        hymn_number = meta.get('hymn_number', '')
        verses = meta.get('verses', [])

        # Intro slide: 4_Section Header with hymn name and number
        slide = self.add_slide('4_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)

        ph1 = self._get_placeholder(slide, 1)
        if ph1:
            tf = ph1.text_frame
            tf.clear()
            # Hymn title (italic)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = hymn_title
            run.font.italic = True
            run.font.size = SUBTITLE_SIZE
            run.font.bold = False
            # Hymn number
            if hymn_number:
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = hymn_number
                run2.font.size = SUBTITLE_SIZE
                run2.font.color.rgb = HYMN_NUM_GOLD

        # One slide per verse: 1_Title and Content
        for verse_num, verse_text in verses:
            slide = self.add_slide('1_Title and Content')
            ph1 = self._get_placeholder(slide, 1)
            if ph1:
                tf = ph1.text_frame
                tf.clear()
                lines = verse_text.split('\n')
                for li, line in enumerate(lines):
                    if li == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = line
                    run.font.italic = True
                    run.font.size = HYMN_LYRIC_SIZE

            # Add verse number text box
            self._add_textbox(slide, str(verse_num),
                              left=Inches(8.5), top=Inches(5.5),
                              width=Inches(1), height=Inches(0.8),
                              size=VERSE_NUM_SIZE, bold=True)

    # ── Greeting ────────────────────────────────────────────────────────

    def _gen_greeting(self, section):
        """Generate the Greeting slide."""
        slide = self.add_slide('1_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)
        content = '\n'.join(section.get('content', []))
        if content:
            self._format_dialogue_in_placeholder(slide, 1, content, size=BODY_SIZE)

    # ── Kyrie ───────────────────────────────────────────────────────────

    def _gen_kyrie(self, section):
        """Generate Kyrie slide(s) - placeholder for hymnal image."""
        slide = self.add_slide('4_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)

        # Add page reference if found
        content = section.get('content', [])
        page_ref = section.get('metadata', {}).get('page_ref', '')
        if not page_ref:
            for line in content:
                match = re.search(r'p\.?\s*\d+', line)
                if match:
                    page_ref = match.group(0)
                    break

        if page_ref:
            self._add_textbox(slide, page_ref,
                              left=Inches(3.5), top=Inches(3),
                              width=Inches(3), height=Inches(1),
                              size=Pt(48))

    # ── Prayer of the Day ───────────────────────────────────────────────

    def _gen_prayer_of_day(self, section):
        """Generate Prayer of the Day slides (1-2 slides)."""
        content = '\n'.join(section.get('content', []))
        chunks = split_text_into_slides(content, max_chars=400, max_lines=7)

        if not chunks:
            slide = self.add_slide('6_Section Header')
            self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)
            return

        for i, chunk in enumerate(chunks):
            slide = self.add_slide('6_Section Header')
            self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)
            # Use placeholder 10 (the lower body text area)
            self._format_dialogue_in_placeholder(slide, 10, chunk, size=BODY_SIZE)

    # ── Children's Message ──────────────────────────────────────────────

    def _gen_childrens_message(self, section):
        """Generate Children's Message slide."""
        slide = self.add_slide('1_Section Header')
        self._set_placeholder_text(slide, 0, "Children's Message", size=TITLE_SIZE)
        content = section.get('content', [])
        if content:
            # Usually just the leader's name
            self._set_placeholder_text(slide, 1, '\n'.join(content), size=BODY_SIZE)

    # ── Readings (First, Second, Psalm) ─────────────────────────────────

    def _gen_reading(self, section):
        """Generate reading slides: intro + text slides."""
        meta = section.get('metadata', {})
        reference = meta.get('reference', '')
        pew_bible = meta.get('pew_bible', '')
        reading_text = meta.get('reading_text', '')

        # Determine display title
        title = section['title']
        if title.startswith('+'):
            title = title[1:].strip()

        # Intro slide: 24_Section Header
        slide = self.add_slide('24_Section Header')
        self._set_placeholder_text(slide, 0, title, size=TITLE_SIZE)

        # Reference and pew bible in body placeholder
        ph1 = self._get_placeholder(slide, 1)
        if ph1:
            tf = ph1.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            # Reference (with spacing)
            ref_display = reference
            if pew_bible:
                # Add spacing between reference and pew bible
                ref_display = reference + '    ' + pew_bible
                run1 = p.add_run()
                run1.text = reference + '    '
                run2 = p.add_run()
                run2.text = pew_bible
                run2.font.bold = False
            else:
                run1 = p.add_run()
                run1.text = reference

        # Reading text slides: 5_Title and Content
        if reading_text:
            chunks = split_reading_text(reading_text)
            for ci, chunk in enumerate(chunks):
                slide = self.add_slide('5_Title and Content')
                ph1 = self._get_placeholder(slide, 1)
                if ph1:
                    tf = ph1.text_frame
                    tf.clear()
                    text = chunk

                    # Check if this is the last chunk - might need "Word of God" response
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = text
                    run.font.size = READING_SIZE
                    run.font.bold = False

    # ── Gospel sections ─────────────────────────────────────────────────

    def _gen_gospel_acclamation(self, section):
        """Generate Gospel Acclamation slide."""
        slide = self.add_slide('6_Section Header')
        # Use a non-placeholder title since the original used a custom text box
        self._set_placeholder_text(slide, 0, '+ Gospel Acclamation', size=Pt(56))

    def _gen_gospel_announcement(self, section):
        """Generate Gospel Announcement slide."""
        slide = self.add_slide('6_Section Header')
        self._set_placeholder_text(slide, 0, '+ Gospel Announcement', size=TITLE_SIZE)

        content = '\n'.join(section.get('content', []))
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_SIZE)

    def _gen_gospel_reading(self, section):
        """Generate Gospel Reading slides."""
        meta = section.get('metadata', {})
        reference = meta.get('reference', '')
        pew_bible = meta.get('pew_bible', '')
        reading_text = meta.get('reading_text', '')
        content = section.get('content', [])

        # Intro slide: 6_Section Header with reference and first part of text
        slide = self.add_slide('6_Section Header')
        self._set_placeholder_text(slide, 0, '+ Gospel Reading', size=TITLE_SIZE)

        # Reference in placeholder 1
        ph1 = self._get_placeholder(slide, 1)
        if ph1:
            tf = ph1.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = reference
            run.font.size = PAGE_REF_SIZE
            if pew_bible:
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = pew_bible
                run2.font.size = PAGE_REF_SIZE
                run2.font.bold = False

        # First portion of reading in placeholder 10
        if reading_text:
            chunks = split_reading_text(reading_text)
            if chunks:
                # Put first chunk in the intro slide's placeholder 10
                self._set_placeholder_text(slide, 10, chunks[0],
                                           size=BODY_SMALL)

                # Remaining chunks on 5_Title and Content slides
                for chunk in chunks[1:]:
                    text_slide = self.add_slide('5_Title and Content')
                    ph = self._get_placeholder(text_slide, 1)
                    if ph:
                        tf = ph.text_frame
                        tf.clear()
                        p = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = chunk
                        run.font.size = READING_SIZE

    # ── Message/Sermon ──────────────────────────────────────────────────

    def _gen_message(self, section):
        """Generate Message slide."""
        slide = self.add_slide('24_Section Header')
        self._set_placeholder_text(slide, 0, 'Message', size=TITLE_SIZE)
        content = section.get('content', [])
        if content:
            self._set_placeholder_text(slide, 1, '\n'.join(content), size=BODY_SIZE)

    # ── Apostles' Creed ─────────────────────────────────────────────────

    def _gen_creed(self, section):
        """Generate Creed slides (~3 slides of bold text)."""
        page_ref = section.get('metadata', {}).get('page_ref', 'p. 105')
        # Filter out page reference lines from content
        content_lines = [l for l in section.get('content', [])
                         if not re.match(r'^p\.?\s*\d+$', l.strip())]
        content = '\n'.join(content_lines)

        # Split creed text across slides
        chunks = split_text_into_slides(content, max_chars=350, max_lines=7)

        for i, chunk in enumerate(chunks):
            slide = self.add_slide('14_Section Header')
            if i == 0:
                self._set_placeholder_text(slide, 0, section['title'],
                                           size=SUBTITLE_SIZE)
                # Page reference in placeholder 1
                if page_ref:
                    self._set_placeholder_text(slide, 1, page_ref,
                                               size=PAGE_REF_SIZE, italic=False)

            # Creed text in placeholder 10 (all bold)
            ph10 = self._get_placeholder(slide, 10)
            if ph10:
                tf = ph10.text_frame
                tf.clear()
                lines = chunk.split('\n')
                for li, line in enumerate(lines):
                    if li == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = line
                    run.font.bold = True
                    run.font.size = BODY_SIZE

    # ── Prayers of Intercession ─────────────────────────────────────────

    def _gen_prayers(self, section):
        """Generate Prayers of Intercession slide."""
        slide = self.add_slide('6_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)

        content = section.get('content', [])
        if content:
            # Usually just the response: "Hear us, O God. Your mercy is great."
            self._format_dialogue_in_placeholder(slide, 10, '\n'.join(content),
                                                  size=BODY_SIZE)

    # ── Sharing of the Peace ────────────────────────────────────────────

    def _gen_peace(self, section):
        """Generate Sharing of the Peace slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, 'Sharing of the Peace', size=TITLE_SIZE)

        page_ref = section.get('metadata', {}).get('page_ref', 'p. 155')
        if page_ref:
            self._set_placeholder_text(slide, 1, page_ref,
                                       size=PAGE_REF_SIZE, italic=False)

        content_lines = [l for l in section.get('content', [])
                         if not re.match(r'^p\.?\s*\d+$', l.strip())]
        content = '\n'.join(content_lines)
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_LARGE)

    # ── Offering ────────────────────────────────────────────────────────

    def _gen_offering(self, section):
        """Generate Offering slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, 'Offering', size=TITLE_SIZE)

        content = section.get('content', [])
        if content:
            # First line usually says what music (e.g., "Choir")
            self._set_placeholder_text(slide, 1, content[0], size=BODY_SIZE,
                                       italic=False)

        # Add offering info text box
        offering_info = [
            'Online offering is available at www.stlukememphis.org or mail to:',
            'St. Luke Lutheran Church',
            '2000 N. Germantown Pkwy.',
        ]
        self._add_textbox(slide, '\n'.join(offering_info),
                          left=Inches(0.7), top=Inches(5),
                          width=Inches(8.5), height=Inches(2),
                          size=Pt(30))

    # ── Offertory ───────────────────────────────────────────────────────

    def _gen_offertory(self, section):
        """Generate Offertory hymn slide(s)."""
        meta = section.get('metadata', {})
        hymn_title = meta.get('hymn_title', '')
        hymn_number = meta.get('hymn_number', '')
        verses = meta.get('verses', [])
        content = section.get('content', [])

        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)

        # Show hymn title + number and lyrics in placeholder 10
        display_parts = []
        if hymn_title and hymn_number:
            display_parts.append(f'{hymn_title} {hymn_number}')
        elif hymn_title:
            display_parts.append(hymn_title)

        # If we have verses, show them
        if verses:
            for _, verse_text in verses:
                display_parts.append(verse_text)
        elif content:
            # Use raw content for lyrics
            # Skip title/number lines
            lyrics_start = 0
            for idx, line in enumerate(content):
                if '#' in line or line == hymn_title:
                    lyrics_start = idx + 1
                else:
                    break
            display_parts.extend(content[lyrics_start:])

        if display_parts:
            ph10 = self._get_placeholder(slide, 10)
            if ph10:
                tf = ph10.text_frame
                tf.clear()
                all_text = '\n'.join(display_parts)
                lines = all_text.split('\n')
                for li, line in enumerate(lines):
                    if li == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = line
                    run.font.size = BODY_SIZE

    # ── Offertory Prayer ────────────────────────────────────────────────

    def _gen_offertory_prayer(self, section):
        """Generate Offertory Prayer slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, '+ Offertory Prayer', size=TITLE_SIZE)

        content = '\n'.join(section.get('content', []))
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_SMALL)

    # ── Great Thanksgiving ──────────────────────────────────────────────

    def _gen_great_thanksgiving(self, section):
        """Generate Great Thanksgiving slides."""
        content = section.get('content', [])
        full_text = '\n'.join(content)

        # Split into chunks
        chunks = split_text_into_slides(full_text, max_chars=400, max_lines=8)

        if not chunks:
            slide = self.add_slide('14_Section Header')
            self._set_placeholder_text(slide, 0, '+ The Great Thanksgiving',
                                       size=SUBTITLE_SIZE)
            return

        # First slide: 14_Section Header with dialogue
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, '+ The Great Thanksgiving',
                                   size=SUBTITLE_SIZE)
        self._format_dialogue_in_placeholder(slide, 10, chunks[0], size=Pt(32))

        # Remaining slides: 5_Title and Content
        for chunk in chunks[1:]:
            slide = self.add_slide('5_Title and Content')
            self._format_dialogue_in_placeholder(slide, 1, chunk, size=BODY_SIZE)

    # ── Words of Institution ────────────────────────────────────────────

    def _gen_words_of_institution(self, section):
        """Generate Words of Institution slides."""
        content = '\n'.join(section.get('content', []))
        chunks = split_text_into_slides(content, max_chars=400, max_lines=7)

        for chunk in chunks:
            slide = self.add_slide('5_Title and Content')
            self._format_dialogue_in_placeholder(slide, 1, chunk, size=BODY_SIZE)

    # ── Lord's Prayer ───────────────────────────────────────────────────

    def _gen_lords_prayer(self, section):
        """Generate Lord's Prayer slides."""
        # Filter out page reference lines from content
        content = [l for l in section.get('content', [])
                   if not re.match(r'^p\.?\s*\d+$', l.strip())]
        page_ref = section.get('metadata', {}).get('page_ref', 'p. 134')

        # Separate intro (P: line) from prayer text
        intro_lines = []
        prayer_lines = []
        in_prayer = False

        for line in content:
            if line.strip().startswith('C:') or in_prayer:
                in_prayer = True
                prayer_lines.append(line.replace('C:', '').strip() if line.strip().startswith('C:') else line)
            else:
                intro_lines.append(line)

        # Intro slide: 14_Section Header
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, "+ The Lord's Prayer", size=TITLE_SIZE)
        if page_ref:
            self._set_placeholder_text(slide, 1, page_ref,
                                       size=PAGE_REF_SIZE, italic=False)
        if intro_lines:
            self._format_dialogue_in_placeholder(slide, 10, '\n'.join(intro_lines),
                                                  size=BODY_LARGE)

        # Prayer text slides: 23_Title and Content
        if prayer_lines:
            prayer_text = '\n'.join(prayer_lines)
            chunks = split_text_into_slides(prayer_text, max_chars=300, max_lines=5)
            for i, chunk in enumerate(chunks):
                slide = self.add_slide('23_Title and Content')
                ph1 = self._get_placeholder(slide, 1)
                if ph1:
                    tf = ph1.text_frame
                    tf.clear()
                    lines = chunk.split('\n')
                    for li, line in enumerate(lines):
                        if li == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.alignment = PP_ALIGN.LEFT
                        run = p.add_run()
                        # Add "C:" prefix on first line of first chunk
                        if i == 0 and li == 0:
                            run.text = 'C: ' + line
                        else:
                            run.text = line

    # ── Communion ───────────────────────────────────────────────────────

    def _gen_communion_invitation(self, section):
        """Generate Invitation to Holy Communion slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, '+ Invitation to Holy Communion',
                                   size=TITLE_SIZE)
        content = '\n'.join(section.get('content', []))
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_LARGE)

    def _gen_communion_hymn(self, section):
        """Generate Communion Hymn intro slide."""
        meta = section.get('metadata', {})
        hymn_title = meta.get('hymn_title', '')
        hymn_number = meta.get('hymn_number', '')

        slide = self.add_slide('4_Section Header')
        self._set_placeholder_text(slide, 0, 'Communion Hymn', size=TITLE_SIZE)

        ph1 = self._get_placeholder(slide, 1)
        if ph1:
            tf = ph1.text_frame
            tf.clear()
            content = section.get('content', [])
            lines_to_show = []
            if hymn_title:
                lines_to_show.append(hymn_title)
            # Add any other content lines (like additional hymn names)
            for line in content:
                if line != hymn_title and line != hymn_number:
                    lines_to_show.append(line)

            for li, line in enumerate(lines_to_show):
                if li == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                # Check for hymn number in line
                num_match = re.search(r'(#\d+)', line)
                if num_match:
                    before = line[:num_match.start()]
                    number = num_match.group(1)
                    after = line[num_match.end():]
                    if before.strip():
                        run = p.add_run()
                        run.text = before
                        run.font.italic = True
                        run.font.size = SUBTITLE_SIZE
                        run.font.bold = False
                    run_num = p.add_run()
                    run_num.text = number
                    run_num.font.size = SUBTITLE_SIZE
                    run_num.font.color.rgb = HYMN_NUM_GOLD
                    if after.strip():
                        run_after = p.add_run()
                        run_after.text = after
                        run_after.font.size = SUBTITLE_SIZE
                else:
                    run = p.add_run()
                    run.text = line
                    run.font.italic = True
                    run.font.size = SUBTITLE_SIZE
                    run.font.bold = False

    # ── Post-Communion ──────────────────────────────────────────────────

    def _gen_post_communion_blessing(self, section):
        """Generate Post Communion Blessing slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, '+ Post Communion Blessing', size=TITLE_SIZE)
        content = '\n'.join(section.get('content', []))
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_SIZE)

    def _gen_post_communion_canticle(self, section):
        """Generate Post Communion Canticle slides."""
        meta = section.get('metadata', {})
        hymn_title = meta.get('hymn_title', '')
        hymn_number = meta.get('hymn_number', '')
        verses = meta.get('verses', [])

        # Intro slide
        slide = self.add_slide('4_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)

        ph1 = self._get_placeholder(slide, 1)
        if ph1:
            tf = ph1.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            if hymn_title:
                run = p.add_run()
                run.text = hymn_title
                run.font.italic = True
                run.font.size = SUBTITLE_SIZE
                run.font.bold = False
            if hymn_number:
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = f'( {hymn_number}, Verse 1 only)'
                run2.font.size = SUBTITLE_SIZE

        # Verse slides
        for verse_num, verse_text in verses:
            slide = self.add_slide('1_Title and Content')
            ph = self._get_placeholder(slide, 1)
            if ph:
                tf = ph.text_frame
                tf.clear()
                for li, line in enumerate(verse_text.split('\n')):
                    if li == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = line
                    run.font.italic = True
                    run.font.size = HYMN_LYRIC_SIZE
            # Verse number
            self._add_textbox(slide, str(verse_num),
                              left=Inches(8.5), top=Inches(5.5),
                              width=Inches(1), height=Inches(0.8),
                              size=VERSE_NUM_SIZE, bold=True)

    def _gen_post_communion_prayer(self, section):
        """Generate Post Communion Prayer slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, '+ Post Communion Prayer', size=TITLE_SIZE)
        content = '\n'.join(section.get('content', []))
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_SMALL)

    # ── Blessing ────────────────────────────────────────────────────────

    def _gen_blessing(self, section):
        """Generate Blessing slides (2 slides)."""
        content = '\n'.join(section.get('content', []))
        chunks = split_text_into_slides(content, max_chars=350, max_lines=5)

        for chunk in chunks:
            slide = self.add_slide('14_Section Header')
            self._set_placeholder_text(slide, 0, '+ Blessing', size=TITLE_SIZE)
            self._format_dialogue_in_placeholder(slide, 10, chunk, size=BODY_SIZE)

    # ── Dismissal ───────────────────────────────────────────────────────

    def _gen_dismissal(self, section):
        """Generate Dismissal slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, '+ Dismissal', size=TITLE_SIZE)

        page_ref = section.get('metadata', {}).get('page_ref', 'p. 155')
        if page_ref:
            self._set_placeholder_text(slide, 1, page_ref,
                                       size=PAGE_REF_SIZE, italic=False)

        content_lines = [l for l in section.get('content', [])
                         if not re.match(r'^p\.?\s*\d+$', l.strip())]
        content = '\n'.join(content_lines)
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_LARGE)

    # ── Postlude ────────────────────────────────────────────────────────

    def _gen_postlude(self, section):
        """Generate Postlude slide."""
        slide = self.add_slide('23_Section Header')
        self._set_placeholder_text(slide, 0, '+Postlude', size=TITLE_SIZE)

        content = section.get('content', [])
        if content:
            # Performer name in placeholder 1
            self._set_placeholder_text(slide, 1, content[0], size=BODY_SIZE)

        # License info and tagline in placeholder 10
        ph10 = self._get_placeholder(slide, 10)
        if ph10:
            tf = ph10.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = 'Seek God. Serve Others. Share Life.'
            run.font.color.rgb = PEACH

    # ── Media placeholder ───────────────────────────────────────────────

    def _gen_media_placeholder(self, section):
        """Generate a placeholder slide for media content (Kyrie, Lamb of God, etc.)."""
        slide = self.add_slide('4_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)

    # ── Generic section ─────────────────────────────────────────────────

    def _gen_generic_section(self, section):
        """Generate a generic section slide."""
        slide = self.add_slide('14_Section Header')
        self._set_placeholder_text(slide, 0, section['title'], size=TITLE_SIZE)
        content = '\n'.join(section.get('content', []))
        if content:
            self._format_dialogue_in_placeholder(slide, 10, content, size=BODY_SIZE)

    # ── Helper methods ──────────────────────────────────────────────────

    def _get_placeholder(self, slide, idx):
        """Get a placeholder by index from a slide."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _set_text(self, placeholder, text, size=None, bold=None, italic=None,
                  color=None):
        """Set simple text in a placeholder."""
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        if size:
            run.font.size = size
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if color:
            run.font.color.rgb = color

    def _set_placeholder_text(self, slide, ph_idx, text, size=None,
                               bold=None, italic=None, color=None):
        """Set text in a placeholder by index."""
        ph = self._get_placeholder(slide, ph_idx)
        if ph:
            self._set_text(ph, text, size=size, bold=bold, italic=italic,
                           color=color)

    def _format_dialogue_in_placeholder(self, slide, ph_idx, text, size=None):
        """Format dialogue text (P:/C: exchanges) in a placeholder.

        C: lines get bold text with gold-colored "C:" prefix.
        P:/AM:/PM: lines get normal weight.
        """
        ph = self._get_placeholder(slide, ph_idx)
        if not ph:
            return

        tf = ph.text_frame
        tf.clear()
        lines = text.split('\n')
        font_size = size or BODY_SIZE

        for li, line in enumerate(lines):
            if li == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            stripped = line.strip()
            if not stripped:
                continue

            if is_congregation_line(stripped):
                # C: prefix in gold, rest bold
                prefix = 'C:'
                rest = stripped[2:].strip()
                # Add "C:" run in gold
                run_prefix = p.add_run()
                run_prefix.text = 'C:'
                run_prefix.font.bold = True
                run_prefix.font.size = font_size
                run_prefix.font.color.rgb = GOLD
                # Add rest of text bold
                run_text = p.add_run()
                run_text.text = '\t' + rest if rest else ''
                run_text.font.bold = True
                run_text.font.size = font_size
            elif stripped.startswith(('P:', 'PM:', 'AM:', 'L:')):
                # Minister lines - normal weight
                run = p.add_run()
                run.text = stripped
                run.font.size = font_size
                run.font.bold = False
            else:
                # Regular text
                run = p.add_run()
                run.text = stripped
                run.font.size = font_size

    def _add_textbox(self, slide, text, left, top, width, height,
                     size=None, bold=None, italic=None, color=None,
                     alignment=None):
        """Add a free-floating text box to a slide."""
        from pptx.util import Inches as In
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if alignment:
            p.alignment = alignment
        run = p.add_run()
        run.text = text
        if size:
            run.font.size = size
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return txBox


# ── Standalone entry point ──────────────────────────────────────────────────

if __name__ == '__main__':
    import sys
    from docx_parser import parse_worship_script

    if len(sys.argv) < 3:
        print("Usage: python slide_generator.py <input.docx> <output.pptx>")
        sys.exit(1)

    sections = parse_worship_script(sys.argv[1])
    gen = SlideGenerator()
    gen.generate(sections)
    gen.save(sys.argv[2])
    print(f"Generated {len(gen.prs.slides)} slides -> {sys.argv[2]}")
